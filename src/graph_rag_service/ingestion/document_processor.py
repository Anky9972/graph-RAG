"""
Document processing and chunking
Gap #9: Extended formats — CSV, Excel, PPTX, JSON (+ existing PDF, TXT, MD, DOCX)
Uses LlamaParse for advanced PDF parsing when available
"""

import aiofiles
from pathlib import Path
from typing import List, Dict, Any, Optional
import hashlib
from datetime import datetime
import os
import json

from pypdf import PdfReader
from llama_index.core.node_parser import SentenceSplitter

from ..core.models import Document, Chunk
from ..config import settings

# LlamaParse import (optional)
try:
    from llama_parse import LlamaParse
    LLAMA_PARSE_AVAILABLE = True
except ImportError:
    LLAMA_PARSE_AVAILABLE = False


class DocumentProcessor:
    """
    Process and chunk documents for ingestion.
    Supports: PDF, TXT, MD, DOCX, CSV, XLSX, PPTX, JSON
    """

    def __init__(self):
        self.chunk_size = settings.chunk_size
        self.chunk_overlap = settings.chunk_overlap
        self.splitter = SentenceSplitter(
            chunk_size=self.chunk_size,
            chunk_overlap=self.chunk_overlap
        )

        # Initialize LlamaParse if API key is available
        self.llama_parser: Optional[LlamaParse] = None
        if (LLAMA_PARSE_AVAILABLE and
                settings.use_llama_parse and
                settings.llama_cloud_api_key):
            try:
                os.environ["LLAMA_CLOUD_API_KEY"] = settings.llama_cloud_api_key
                self.llama_parser = LlamaParse(
                    result_type="markdown",
                    verbose=True,
                    language="en",
                )
            except Exception as e:
                print(f"Warning: Failed to initialize LlamaParse: {e}")
                self.llama_parser = None

    async def process_document(self, file_path: Path) -> Document:
        """
        Process a document and extract metadata

        Args:
            file_path: Path to document file

        Returns:
            Document with metadata
        """
        text = await self._extract_text(file_path)

        document = Document(
            id=self._generate_document_id(file_path),
            filename=file_path.name,
            file_type=file_path.suffix,
            size_bytes=file_path.stat().st_size,
            upload_date=datetime.utcnow(),
            content=text,
            metadata={
                "file_path": str(file_path),
                "extension": file_path.suffix
            }
        )

        return document

    async def chunk_document(
        self,
        document: Document
    ) -> List[Chunk]:
        """
        Chunk document into smaller pieces

        Args:
            document: Document to chunk

        Returns:
            List of chunks
        """
        if not document.content:
            return []

        text_chunks = self.splitter.split_text(document.content)

        chunks = []
        for i, text in enumerate(text_chunks):
            chunk = Chunk(
                id=f"{document.id}_chunk_{i}",
                text=text,
                document_id=document.id,
                chunk_index=i,
                metadata={
                    "document_filename": document.filename,
                    "document_type": document.file_type,
                    "chunk_index": i,
                    "total_chunks": len(text_chunks)
                }
            )
            chunks.append(chunk)

        return chunks

    async def _extract_text(self, file_path: Path) -> str:
        """Extract text from file based on type"""
        extension = file_path.suffix.lower()

        if extension == '.pdf':
            return await self._extract_pdf(file_path)
        elif extension in ('.txt', '.md'):
            return await self._extract_txt(file_path)
        elif extension == '.docx':
            return await self._extract_docx(file_path)
        # ── Gap #9: New formats ───────────────────────────────────────────────
        elif extension == '.csv':
            return await self._extract_csv(file_path)
        elif extension in ('.xlsx', '.xls'):
            return await self._extract_excel(file_path)
        elif extension == '.pptx':
            return await self._extract_pptx(file_path)
        elif extension == '.json':
            return await self._extract_json(file_path)
        else:
            raise ValueError(f"Unsupported file type: {extension}")

    # ── Existing extractors ───────────────────────────────────────────────────

    async def _extract_pdf(self, file_path: Path) -> str:
        """Extract text from PDF using LlamaParse (if available) or pypdf"""
        if self.llama_parser:
            try:
                print(f"Using LlamaParse for {file_path.name}...")
                documents = await self.llama_parser.aload_data(str(file_path))
                text = "\n\n".join([doc.text for doc in documents])
                print(f"✓ LlamaParse extracted {len(text)} characters")
                return text.strip()
            except Exception as e:
                print(f"Warning: LlamaParse failed, falling back to pypdf: {e}")

        print(f"Using pypdf for {file_path.name}...")
        reader = PdfReader(str(file_path))
        text = ""
        for page_num, page in enumerate(reader.pages):
            page_text = page.extract_text()
            if page_text:
                text += f"\n[Page {page_num + 1}]\n{page_text}\n"
        return text.strip()

    async def _extract_txt(self, file_path: Path) -> str:
        """Extract text from TXT/MD file"""
        async with aiofiles.open(file_path, 'r', encoding='utf-8') as f:
            return await f.read()

    async def _extract_docx(self, file_path: Path) -> str:
        """Extract text from DOCX"""
        try:
            import zipfile
            import xml.etree.ElementTree as ET

            with zipfile.ZipFile(file_path) as docx:
                xml_content = docx.read('word/document.xml')
                tree = ET.XML(xml_content)

                paragraphs = []
                for paragraph in tree.iter():
                    if paragraph.tag.endswith('}t'):
                        if paragraph.text:
                            paragraphs.append(paragraph.text)

                return '\n'.join(paragraphs)
        except Exception as e:
            raise ValueError(f"Failed to extract text from DOCX: {e}")

    # ── Gap #9: New format extractors ─────────────────────────────────────────

    async def _extract_csv(self, file_path: Path) -> str:
        """
        Extract CSV as structured text.
        Each row becomes a natural language sentence-like string.
        This allows the LLM extractor to identify entities from tabular data.
        """
        try:
            import csv
            lines = []
            async with aiofiles.open(file_path, 'r', encoding='utf-8', newline='') as f:
                content = await f.read()

            reader = csv.DictReader(content.splitlines())
            headers = reader.fieldnames or []

            lines.append(f"CSV Data from: {file_path.name}")
            lines.append(f"Columns: {', '.join(headers)}")
            lines.append("")

            for i, row in enumerate(reader):
                # Convert each row to a descriptive sentence
                parts = [f"{k}={v}" for k, v in row.items() if v and v.strip()]
                lines.append(f"Row {i+1}: " + " | ".join(parts))

            return "\n".join(lines)
        except Exception as e:
            raise ValueError(f"Failed to extract CSV: {e}")

    async def _extract_excel(self, file_path: Path) -> str:
        """
        Extract Excel spreadsheet content.
        Processes all sheets, converts to structured text.
        """
        try:
            import openpyxl
            wb = openpyxl.load_workbook(str(file_path), data_only=True)
            all_text = []

            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                all_text.append(f"\n=== Sheet: {sheet_name} ===\n")

                # Get headers from first row
                headers = []
                first_row = True
                for row in ws.iter_rows(values_only=True):
                    if all(v is None for v in row):
                        continue
                    if first_row:
                        headers = [str(v) if v is not None else "" for v in row]
                        all_text.append(f"Columns: {', '.join(h for h in headers if h)}")
                        first_row = False
                        continue
                    # Format each data row
                    parts = []
                    for header, value in zip(headers, row):
                        if value is not None and str(value).strip():
                            parts.append(f"{header}={value}")
                    if parts:
                        all_text.append(" | ".join(parts))

            return "\n".join(all_text)
        except ImportError:
            raise ValueError("openpyxl not installed. Run: pip install openpyxl")
        except Exception as e:
            raise ValueError(f"Failed to extract Excel: {e}")

    async def _extract_pptx(self, file_path: Path) -> str:
        """
        Extract PowerPoint presentation content.
        Processes each slide: title + body text + speaker notes.
        """
        try:
            from pptx import Presentation
            prs = Presentation(str(file_path))
            slides_text = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_parts = [f"\n=== Slide {slide_num} ==="]

                # Title
                if slide.shapes.title and slide.shapes.title.text:
                    slide_parts.append(f"Title: {slide.shapes.title.text.strip()}")

                # Body text
                body_texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame and shape != slide.shapes.title:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                body_texts.append(text)
                if body_texts:
                    slide_parts.append("Content:\n" + "\n".join(body_texts))

                # Speaker notes
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        slide_parts.append(f"Notes: {notes_text}")

                slides_text.append("\n".join(slide_parts))

            return "\n\n".join(slides_text)
        except ImportError:
            raise ValueError("python-pptx not installed. Run: pip install python-pptx")
        except Exception as e:
            raise ValueError(f"Failed to extract PPTX: {e}")

    async def _extract_json(self, file_path: Path) -> str:
        """
        Extract JSON content.
        Flattens nested structures into readable text for entity extraction.
        """
        try:
            async with aiofiles.open(file_path, 'r', encoding='utf-8') as f:
                content = await f.read()

            data = json.loads(content)
            lines = [f"JSON Data from: {file_path.name}", ""]

            def flatten(obj: Any, prefix: str = "") -> List[str]:
                parts = []
                if isinstance(obj, dict):
                    for k, v in obj.items():
                        key = f"{prefix}.{k}" if prefix else k
                        parts.extend(flatten(v, prefix=key))
                elif isinstance(obj, list):
                    for i, item in enumerate(obj[:50]):  # limit list items
                        key = f"{prefix}[{i}]"
                        parts.extend(flatten(item, prefix=key))
                else:
                    if obj is not None and str(obj).strip():
                        parts.append(f"{prefix}: {obj}")
                return parts

            if isinstance(data, list):
                lines.append(f"Array with {len(data)} items:")
                for i, item in enumerate(data[:100]):  # limit root array
                    lines.append(f"\nItem {i+1}:")
                    lines.extend(flatten(item))
            else:
                lines.extend(flatten(data))

            return "\n".join(lines)
        except Exception as e:
            raise ValueError(f"Failed to extract JSON: {e}")

    def _generate_document_id(self, file_path: Path) -> str:
        """Generate unique document ID based on file content"""
        hasher = hashlib.sha256()
        hasher.update(str(file_path).encode())
        hasher.update(str(file_path.stat().st_mtime).encode())
        return hasher.hexdigest()[:16]
